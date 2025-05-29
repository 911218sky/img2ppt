import os
from typing import List, Tuple

from natsort import natsorted
from PIL import Image
from pptx import Presentation
from pptx.util import Emu

from config import (
    EMU_PER_PX,
    SLIDE_WIDTH,
    SLIDE_HEIGHT,
    IMAGE_EXTENSIONS
)

def get_image_files(
    image_dir: str,
    exts: Tuple[str, ...] = IMAGE_EXTENSIONS
) -> List[str]:
    """
    Retrieve all image files in a directory and return them
    sorted in natural (human) numeric order.

    :param image_dir: Path to the folder containing images
    :param exts: Tuple of supported image file extensions
    :return: Naturally sorted list of image filenames
    """
    files = [
        fname for fname in os.listdir(image_dir)
        if fname.lower().endswith(exts)
    ]
    return natsorted(files)

def images_to_pptx_keep_original(
    image_dir: str,
    output_pptx: str = 'output.pptx'
) -> None:
    """
    Create a 16:9 PowerPoint presentation where each slide
    contains one image at its original resolution. Images
    larger than the slide will be scaled down proportionally
    and centered.

    :param image_dir: Folder containing image files
    :param output_pptx: Name of the resulting PPTX file
    """
    # Create a new presentation
    prs = Presentation()
    # Set slide size to 16:9
    prs.slide_width  = Emu(SLIDE_WIDTH)
    prs.slide_height = Emu(SLIDE_HEIGHT)

    image_files = get_image_files(image_dir)
    print(f"Found {len(image_files)} images:")
    print(image_files)

    for img_name in image_files:
        img_path = os.path.join(image_dir, img_name)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        
        # Open image to get its pixel dimensions
        with Image.open(img_path) as im:
            px_w, px_h = im.size
        
        # Convert pixels to EMU (PowerPoint unit)
        pic_w = Emu(px_w * EMU_PER_PX)
        pic_h = Emu(px_h * EMU_PER_PX)

        # Calculate scale factor to fit within slide
        sw, sh = prs.slide_width, prs.slide_height
        scale = min(1.0, sw / pic_w, sh / pic_h)
        pic_w *= scale
        pic_h *= scale

        # Center the image on the slide
        left = (sw - pic_w) / 2
        top  = (sh - pic_h) / 2

        # Add the picture to the slide
        slide.shapes.add_picture(img_path, left, top, width=pic_w, height=pic_h)

    # Save the presentation
    prs.save(output_pptx)
    print(f"PPTX file created: {output_pptx}")

if __name__ == '__main__':
    import sys

    # Prompt user for input folder and output filename
    image_dir = input("Enter the path to the image folder: ").strip()
    if not image_dir:
        print("Error: No image folder provided.")
        sys.exit(1)

    default_name = os.path.basename(os.path.normpath(image_dir)) + '.pptx'
    output_pptx = input(f"Enter output PPTX filename [{default_name}]: ").strip()
    if not output_pptx:
        output_pptx = default_name

    # Run the conversion
    images_to_pptx_keep_original(image_dir, output_pptx)